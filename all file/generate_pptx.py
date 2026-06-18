"""
SmartInvoice AI — PFE Presentation Generator
Generates a high-impact, visually stunning PowerPoint with Morph transitions.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

# ── Colors ──────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x25, 0x40)
BLUE      = RGBColor(0x2E, 0x75, 0xB6)
EMERALD   = RGBColor(0x10, 0xB9, 0x81)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
GRAY      = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_NAVY = RGBColor(0x0E, 0x2F, 0x50)
ACCENT_BLUE = RGBColor(0x38, 0x8A, 0xE5)
SOFT_GREEN = RGBColor(0x06, 0xD6, 0xA0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Utility helpers ─────────────────────────────────────────────────────
def set_morph(slide):
    """Apply Morph transition to slide via raw XML."""
    from lxml import etree

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
    P159_NS = 'http://schemas.microsoft.com/office/powerpoint/2015/09/main'

    nsmap = {'p': P_NS, 'p14': P14_NS, 'p159': P159_NS}

    transition = etree.SubElement(
        slide._element, f'{{{P_NS}}}transition',
        attrib={'spd': 'slow', 'advClick': '1'},
        nsmap=nsmap
    )
    etree.SubElement(
        transition,
        f'{{{P159_NS}}}morph',
        attrib={f'{{{P159_NS}}}option': 'byObject'},
        nsmap=nsmap
    )


def add_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, name=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    if name:
        shape.name = name
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, name=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if name:
        shape.name = name
    return shape


def add_circle(slide, left, top, size, fill_color, name=None):
    """Add a circle (oval)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if name:
        shape.name = name
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI", name=None,
             anchor=MSO_ANCHOR.TOP):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    if name:
        txBox.name = name
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=WHITE, bold=False, spacing=1.2, font_name="Segoe UI",
                  name=None, alignment=PP_ALIGN.LEFT, bullet=False):
    """Add multi-line text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(font_size * spacing * 0.3)
        run = p.add_run()
        if bullet:
            run.text = f"  {line}"
        else:
            run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    if name:
        txBox.name = name
    return txBox


def add_card(slide, left, top, width, height, title, subtitle, icon_text="",
             bg_color=LIGHT_NAVY, title_color=EMERALD, text_color=WHITE, name=None):
    """Add a styled card."""
    card = add_rounded_rect(slide, left, top, width, height, bg_color, name=name)

    if icon_text:
        add_text(slide, left + Inches(0.3), top + Inches(0.25), width - Inches(0.6), Inches(0.5),
                 icon_text, font_size=28, color=EMERALD, bold=True, font_name="Segoe UI Emoji")

    add_text(slide, left + Inches(0.3), top + Inches(0.85), width - Inches(0.6), Inches(0.6),
             title, font_size=14, color=title_color, bold=True)

    add_text(slide, left + Inches(0.3), top + Inches(1.25), width - Inches(0.6), Inches(0.5),
             subtitle, font_size=11, color=GRAY)

    return card


def add_metric_card(slide, left, top, width, height, value, label, accent=EMERALD, name=None):
    """Add a metric card with big number."""
    card = add_rounded_rect(slide, left, top, width, height, LIGHT_NAVY, name=name)

    # Accent top bar
    add_rect(slide, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Pt(3), accent)

    add_text(slide, left + Inches(0.15), top + Inches(0.35), width - Inches(0.3), Inches(0.7),
             value, font_size=32, color=WHITE, bold=True)

    add_text(slide, left + Inches(0.15), top + Inches(1.05), width - Inches(0.3), Inches(0.4),
             label, font_size=11, color=GRAY)

    return card


def add_arrow_flow(slide, items, y, start_x, item_width, gap, colors=None, name_prefix="flow"):
    """Add horizontal flow with arrows."""
    if colors is None:
        colors = [BLUE, EMERALD, ACCENT_BLUE, SOFT_GREEN]
    x = start_x
    for i, item in enumerate(items):
        c = colors[i % len(colors)]
        add_rounded_rect(slide, x, y, item_width, Inches(0.7), c, name=f"{name_prefix}_{i}")
        add_text(slide, x, y + Inches(0.12), item_width, Inches(0.5),
                 item, font_size=12, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Segoe UI")

        if i < len(items) - 1:
            # Arrow
            arrow_x = x + item_width + Inches(0.05)
            add_text(slide, arrow_x, y + Inches(0.08), Inches(0.3), Inches(0.5),
                     "\u2192", font_size=22, color=GRAY, bold=True,
                     alignment=PP_ALIGN.CENTER)

        x += item_width + Inches(0.45)


def add_footer(slide, slide_num, total=22):
    """Add slide footer with number."""
    # Bottom line
    add_rect(slide, Inches(0), H - Inches(0.5), W, Pt(1), EMERALD)
    add_text(slide, W - Inches(1.5), H - Inches(0.45), Inches(1.2), Inches(0.35),
             f"{slide_num} / {total}", font_size=10, color=GRAY,
             alignment=PP_ALIGN.RIGHT, font_name="Segoe UI Light")
    add_text(slide, Inches(0.5), H - Inches(0.45), Inches(3), Inches(0.35),
             "SmartInvoice AI", font_size=10, color=EMERALD,
             font_name="Segoe UI Light", bold=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, NAVY)

# Decorative elements
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(10), Inches(4.5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

# Emerald accent line
add_rect(slide, Inches(1.5), Inches(2.2), Inches(0.8), Pt(4), EMERALD, name="accent_line")

# Title
add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
         "SmartInvoice AI", font_size=56, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="main_title")

# Tagline
add_text(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
         "Plateforme SaaS Multi-Tenant de Gestion Intelligente des Factures",
         font_size=22, color=EMERALD, font_name="Segoe UI Light", name="tagline")

add_text(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.6),
         "avec Conformite Fiscale Tunisienne",
         font_size=20, color=GRAY, font_name="Segoe UI Light", name="tagline2")

# Student info
add_text(slide, Inches(1.5), Inches(5.5), Inches(5), Inches(0.4),
         "Anas Badri  |  Annee academique 2025-2026",
         font_size=14, color=GRAY, font_name="Segoe UI Light", name="student_info")

# Decorative dots
for i in range(5):
    add_circle(slide, Inches(1.5 + i * 0.25), Inches(5.2), Inches(0.08),
               EMERALD if i == 0 else RGBColor(0x1A, 0x3A, 0x55), name=f"dot_{i}")


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problematique
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

# Decorative elements (morph from slide 1)
add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(6), Inches(0.7),
         "La Problematique", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Problem cards
problems = [
    ("Gestion Manuelle", "80% des PME tunisiennes traitent\nleurs factures manuellement", "01"),
    ("Non-Conformite", "Loi El Fatoora 2026 : obligation\nde facturation electronique", "02"),
    ("Erreurs Humaines", "Saisie manuelle = 3-5% taux\nd'erreur sur les montants", "03"),
    ("Temps Perdu", "15-20 min par facture en\ntraitement manuel moyen", "04"),
]

for i, (title, desc, num) in enumerate(problems):
    x = Inches(0.8 + (i % 2) * 5.8)
    y = Inches(2.0 + (i // 2) * 2.5)

    card = add_rounded_rect(slide, x, y, Inches(5.4), Inches(2.1), LIGHT_NAVY, name=f"prob_card_{i}")

    # Number
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.6),
             num, font_size=36, color=EMERALD, bold=True, font_name="Segoe UI Light")

    # Title
    add_text(slide, x + Inches(1.2), y + Inches(0.3), Inches(3.8), Inches(0.5),
             title, font_size=20, color=WHITE, bold=True)

    # Description
    add_text(slide, x + Inches(1.2), y + Inches(0.85), Inches(3.8), Inches(1.0),
             desc.replace('\n', ' '), font_size=13, color=GRAY)

add_footer(slide, 2)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Proposee
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(8), Inches(0.7),
         "Notre Solution", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
         "Une plateforme SaaS complete combinant IA et conformite fiscale tunisienne",
         font_size=16, color=GRAY, font_name="Segoe UI Light")

# Three pillars
pillars = [
    ("Extraction IA", "Pipeline intelligent Claude +\nOllama avec scoring de confiance\net validation automatique", BLUE),
    ("Multi-Tenant SaaS", "Architecture isolee par organisation,\ngestion d'equipe avec roles,\nfacturation Stripe integree", EMERALD),
    ("Conformite TTN", "El Fatoora, TEJ retenue a la source,\nsignature electronique XAdES-B,\nvalidation fiscale 2026", RGBColor(0xF5, 0x9E, 0x0B)),
]

for i, (title, desc, color) in enumerate(pillars):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.8)

    # Card
    add_rounded_rect(slide, x, y, Inches(3.6), Inches(3.8), LIGHT_NAVY, name=f"pillar_{i}")

    # Top accent
    add_rect(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Pt(4), color)

    # Number
    add_text(slide, x + Inches(0.3), y + Inches(0.55), Inches(3), Inches(0.6),
             f"0{i+1}", font_size=42, color=color, bold=True, font_name="Segoe UI Light")

    # Title
    add_text(slide, x + Inches(0.3), y + Inches(1.3), Inches(3), Inches(0.5),
             title, font_size=20, color=WHITE, bold=True)

    # Description
    add_text(slide, x + Inches(0.3), y + Inches(1.9), Inches(3), Inches(1.5),
             desc.replace('\n', ' '), font_size=13, color=GRAY, font_name="Segoe UI")

add_footer(slide, 3)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture Globale
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(8), Inches(0.7),
         "Architecture Multi-Tenant", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Architecture layers
layers = [
    ("FRONTEND", "Next.js 16 / React 19 / TypeScript / Tailwind CSS 4", BLUE, Inches(1.8)),
    ("API GATEWAY", "Flask REST API  |  70+ Endpoints  |  JWT Auth  |  CORS", ACCENT_BLUE, Inches(2.7)),
    ("SERVICES", "Extraction IA  |  TTN  |  TEJ  |  Notifications  |  Billing", EMERALD, Inches(3.6)),
    ("WORKERS", "Celery  |  Redis Broker  |  Traitement Asynchrone", RGBColor(0xF5, 0x9E, 0x0B), Inches(4.5)),
    ("DATA", "MongoDB 7  |  Redis Cache  |  File Storage", RGBColor(0xEF, 0x44, 0x44), Inches(5.4)),
]

for label, desc, color, y in layers:
    # Layer bar
    add_rounded_rect(slide, Inches(1.5), y, Inches(10.3), Inches(0.75), LIGHT_NAVY, name=f"layer_{label}")
    # Color accent
    add_rect(slide, Inches(1.5), y, Inches(0.12), Inches(0.75), color)
    # Label
    add_text(slide, Inches(1.85), y + Inches(0.08), Inches(2.2), Inches(0.35),
             label, font_size=13, color=color, bold=True)
    # Description
    add_text(slide, Inches(4.0), y + Inches(0.08), Inches(7.5), Inches(0.35),
             desc, font_size=12, color=GRAY)

    # Arrow down
    if y < Inches(5.4):
        add_text(slide, Inches(6.5), y + Inches(0.6), Inches(0.5), Inches(0.4),
                 "\u25BC", font_size=14, color=RGBColor(0x3A, 0x50, 0x68),
                 alignment=PP_ALIGN.CENTER)

# Side box: Multi-tenant isolation
add_rounded_rect(slide, Inches(9.5), Inches(1.5), Inches(3.3), Inches(0.9), LIGHT_NAVY)
add_text(slide, Inches(9.7), Inches(1.55), Inches(2.9), Inches(0.3),
         "Isolation Multi-Tenant", font_size=12, color=EMERALD, bold=True)
add_text(slide, Inches(9.7), Inches(1.85), Inches(2.9), Inches(0.4),
         "Donnees isolees par org_id", font_size=10, color=GRAY)

add_footer(slide, 4)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Stack Technique
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(8), Inches(0.7),
         "Stack Technique", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# 4 columns
stacks = [
    ("Frontend", BLUE, [
        "Next.js 16",
        "React 19",
        "TypeScript 5",
        "Tailwind CSS 4",
        "Shadcn/ui + Radix",
        "Recharts 3",
        "React Hook Form",
    ]),
    ("Backend", EMERALD, [
        "Flask + CORS",
        "Python 3.10",
        "PyMongo",
        "Celery + Redis",
        "Gunicorn",
        "pdfplumber",
        "Tesseract OCR",
    ]),
    ("Intelligence Artificielle", RGBColor(0xF5, 0x9E, 0x0B), [
        "Claude Sonnet 4",
        "Ollama (Qwen 2.5)",
        "Template Fallback",
        "Scoring Confiance",
        "OCR Multi-langue",
        "50+ Devises",
        "Validation Auto",
    ]),
    ("Infrastructure", RGBColor(0xEF, 0x44, 0x44), [
        "Docker Compose",
        "MongoDB 7",
        "Redis 7",
        "Stripe API",
        "Twilio SMS",
        "Gmail SMTP",
        "GitHub CI/CD",
    ]),
]

for i, (title, color, items) in enumerate(stacks):
    x = Inches(0.6 + i * 3.1)
    y = Inches(1.9)

    add_rounded_rect(slide, x, y, Inches(2.85), Inches(5.0), LIGHT_NAVY, name=f"stack_{i}")
    add_rect(slide, x, y, Inches(2.85), Inches(0.65), color)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), Inches(2.5), Inches(0.5),
             title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        iy = y + Inches(0.85 + j * 0.55)
        add_text(slide, x + Inches(0.3), iy, Inches(2.3), Inches(0.4),
                 f"\u2022  {item}", font_size=12, color=GRAY)

add_footer(slide, 5)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Pipeline IA Triple Fallback
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.7),
         "Pipeline IA : Triple Fallback", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Flow: Input -> OCR -> LLM -> Output
# Step 1: Input
add_rounded_rect(slide, Inches(0.8), Inches(2.2), Inches(2.5), Inches(1.2), LIGHT_NAVY, name="pipe_input")
add_text(slide, Inches(0.8), Inches(2.3), Inches(2.5), Inches(0.4),
         "ENTREE", font_size=14, color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(2.7), Inches(2.5), Inches(0.5),
         "PDF / Image / XML", font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(3.3), Inches(2.5), Inches(0.5), Inches(0.5),
         "\u2192", font_size=28, color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

# Step 2: OCR
add_rounded_rect(slide, Inches(3.8), Inches(2.2), Inches(2.5), Inches(1.2), LIGHT_NAVY, name="pipe_ocr")
add_text(slide, Inches(3.8), Inches(2.3), Inches(2.5), Inches(0.4),
         "PREPROCESSING", font_size=14, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(3.8), Inches(2.7), Inches(2.5), Inches(0.5),
         "Tesseract OCR + pdfplumber", font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(6.3), Inches(2.5), Inches(0.5), Inches(0.5),
         "\u2192", font_size=28, color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

# Step 3: LLM pipeline
add_rounded_rect(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(1.6), LIGHT_NAVY, name="pipe_llm")
add_text(slide, Inches(6.8), Inches(2.05), Inches(5.8), Inches(0.35),
         "EXTRACTION LLM", font_size=14, color=RGBColor(0xF5, 0x9E, 0x0B),
         bold=True, alignment=PP_ALIGN.CENTER)

# Three LLM boxes inside
llm_items = [
    ("Claude Sonnet 4", "API Cloud\nPrimaire", BLUE),
    ("Ollama Qwen 2.5", "Local\nFallback", EMERALD),
    ("Template", "Dernier\nRecours", GRAY),
]
for j, (name, desc, c) in enumerate(llm_items):
    lx = Inches(7.0 + j * 1.85)
    add_rounded_rect(slide, lx, Inches(2.5), Inches(1.6), Inches(0.9), RGBColor(0x14, 0x38, 0x5C))
    add_text(slide, lx, Inches(2.55), Inches(1.6), Inches(0.3),
             name, font_size=10, color=c, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, lx, Inches(2.85), Inches(1.6), Inches(0.5),
             desc.replace('\n', ' '), font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

# Output section
add_text(slide, Inches(0.8), Inches(4.0), Inches(12), Inches(0.35),
         "\u25BC  SORTIE STRUCTUREE", font_size=14, color=EMERALD, bold=True)

# Output fields
output_cards = [
    ("Donnees Extraites", "vendor, invoice_no, date,\nline_items, totals, bill_to"),
    ("Scoring Confiance", "0-100 par champ,\nchamps critiques identifies"),
    ("Validation Auto", "Conformite fiscale,\nreconciliation arithmetique"),
    ("Review Humain", "Seuil < 60% confiance,\ninterface double ecran"),
]

for i, (title, desc) in enumerate(output_cards):
    x = Inches(0.8 + i * 3.1)
    y = Inches(4.6)
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.2), LIGHT_NAVY)
    add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.4), Inches(0.4),
             title, font_size=14, color=EMERALD, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.4), Inches(1.4),
             desc.replace('\n', ' '), font_size=11, color=GRAY)

add_footer(slide, 6)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Conformite TTN
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.7),
         "Conformite Fiscale Tunisienne", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Three compliance pillars
comp_items = [
    ("El Fatoora TTN", RGBColor(0xF5, 0x9E, 0x0B), [
        "Soumission electronique",
        "Matricule fiscale valide",
        "Format XML structure",
        "Audit trail complet",
        "Score conformite 0-100%",
    ]),
    ("TEJ - Retenue a la Source", BLUE, [
        "Certificats de retenue",
        "5+ types de retenue",
        "Calcul automatique taux",
        "Export PDF signe",
        "Import batch CSV",
    ]),
    ("Signature XAdES-B", EMERALD, [
        "Signature electronique",
        "Validation certificats",
        "Chaine de confiance",
        "NGSign compatible",
        "Verification integree",
    ]),
]

for i, (title, color, items) in enumerate(comp_items):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.0)

    add_rounded_rect(slide, x, y, Inches(3.6), Inches(4.8), LIGHT_NAVY, name=f"comp_{i}")
    add_rect(slide, x, y, Inches(3.6), Inches(0.6), color)
    add_text(slide, x, y + Inches(0.1), Inches(3.6), Inches(0.45),
             title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        iy = y + Inches(0.9 + j * 0.7)
        add_rounded_rect(slide, x + Inches(0.2), iy, Inches(3.2), Inches(0.55),
                         RGBColor(0x14, 0x38, 0x5C))
        add_text(slide, x + Inches(0.4), iy + Inches(0.1), Inches(2.8), Inches(0.35),
                 item, font_size=12, color=GRAY)

add_footer(slide, 7)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Sprint 1 : Authentification
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

# Sprint badge
add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), BLUE, name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 1", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "Authentification & Multi-Tenant", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Features grid
features_s1 = [
    ("JWT HS256", "Token avec claims org_id\net org_role, expiration 24h"),
    ("Inscription", "Onboarding wizard multi-etapes\navec creation org auto"),
    ("Roles RBAC", "Owner, Admin, Member, Viewer\navec decorateurs Python"),
    ("Multi-Tenant", "Isolation donnees par org_id\nsur toutes les collections"),
    ("Invitations", "Systeme d'invitation email\navec token 7 jours"),
    ("Org Switcher", "Navigation entre organisations\nsans re-authentification"),
]

for i, (title, desc) in enumerate(features_s1):
    x = Inches(0.8 + (i % 3) * 4.0)
    y = Inches(2.5 + (i // 3) * 2.3)

    add_rounded_rect(slide, x, y, Inches(3.6), Inches(1.9), LIGHT_NAVY, name=f"s1_feat_{i}")

    add_circle(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.4), BLUE)
    add_text(slide, x + Inches(0.22), y + Inches(0.22), Inches(0.4), Inches(0.35),
             f"0{i+1}", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.8), y + Inches(0.2), Inches(2.6), Inches(0.4),
             title, font_size=16, color=WHITE, bold=True)

    add_text(slide, x + Inches(0.8), y + Inches(0.65), Inches(2.6), Inches(1.0),
             desc.replace('\n', ' '), font_size=11, color=GRAY)

add_footer(slide, 8)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Sprint 1 : Fonctionnalites Cles
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), BLUE, name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 1", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "Flux d'Authentification JWT", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# JWT Flow diagram
flow_steps = [
    ("Client", "Login Request\nemail + password"),
    ("API Flask", "Verification bcrypt\n+ Generation JWT"),
    ("MongoDB", "Lookup user\n+ org membership"),
    ("Token JWT", "user_id, email,\norg_id, org_role, exp"),
    ("Routes Protegees", "@token_required\n@org_required"),
]

for i, (title, desc) in enumerate(flow_steps):
    x = Inches(0.5 + i * 2.5)
    y = Inches(2.5)
    c = [BLUE, EMERALD, RGBColor(0xF5, 0x9E, 0x0B), RGBColor(0xEF, 0x44, 0x44), ACCENT_BLUE][i]

    add_rounded_rect(slide, x, y, Inches(2.2), Inches(1.8), LIGHT_NAVY, name=f"jwt_step_{i}")
    add_rect(slide, x, y, Inches(2.2), Inches(0.45), c)
    add_text(slide, x, y + Inches(0.08), Inches(2.2), Inches(0.3),
             title, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), y + Inches(0.6), Inches(1.9), Inches(1.0),
             desc.replace('\n', ' '), font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(flow_steps) - 1:
        add_text(slide, x + Inches(2.1), y + Inches(0.5), Inches(0.5), Inches(0.5),
                 "\u2192", font_size=22, color=EMERALD, alignment=PP_ALIGN.CENTER)

# Access control table
add_text(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
         "Controle d'Acces par Role", font_size=18, color=EMERALD, bold=True)

roles = [
    ("Owner", "Tout, facturation, supprimer org"),
    ("Admin", "Gerer membres, factures, parametres"),
    ("Member", "Upload, voir, editer factures"),
    ("Viewer", "Lecture seule"),
]
for i, (role, perms) in enumerate(roles):
    y = Inches(5.3 + i * 0.45)
    add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.38), LIGHT_NAVY)
    add_text(slide, Inches(1.0), y + Inches(0.04), Inches(1.5), Inches(0.3),
             role, font_size=12, color=EMERALD, bold=True)
    add_text(slide, Inches(2.8), y + Inches(0.04), Inches(9), Inches(0.3),
             perms, font_size=11, color=GRAY)

add_footer(slide, 9)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Sprint 1 : Metriques
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), BLUE, name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 1", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "Resultats Sprint 1", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Metrics
metrics_s1 = [
    ("4", "Routes Auth", BLUE),
    ("5", "Roles RBAC", EMERALD),
    ("11", "Collections DB", RGBColor(0xF5, 0x9E, 0x0B)),
    ("3", "Contexts React", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (val, label, color) in enumerate(metrics_s1):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.5)
    add_metric_card(slide, x, y, Inches(2.8), Inches(1.6), val, label, color, name=f"s1_metric_{i}")

# Deliverables
add_text(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.4),
         "Livrables Sprint 1", font_size=18, color=EMERALD, bold=True)

deliverables = [
    "Systeme d'inscription + login JWT complet",
    "Architecture multi-tenant avec isolation par org_id",
    "Onboarding wizard multi-etapes fonctionnel",
    "Systeme d'invitations avec tokens securises",
    "Interface org switcher pour navigation multi-org",
    "Decorateurs Python pour controle d'acces (4 niveaux)",
]

for i, d in enumerate(deliverables):
    y = Inches(5.0 + i * 0.38)
    add_text(slide, Inches(1.0), y, Inches(11), Inches(0.35),
             f"\u2713  {d}", font_size=12, color=GRAY)

add_footer(slide, 10)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Sprint 2 : Extraction IA
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), EMERALD, name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 2", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "Extraction IA des Factures", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Pipeline visual
add_text(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.4),
         "Pipeline d'Extraction Intelligente", font_size=16, color=EMERALD, bold=True)

pipeline_steps = [
    ("Upload", "PDF, Image,\nXML", BLUE),
    ("OCR", "Tesseract +\npdfplumber", ACCENT_BLUE),
    ("LLM", "Claude / Ollama\nPrompt 83 lignes", EMERALD),
    ("Validation", "Scoring 0-100\npar champ", RGBColor(0xF5, 0x9E, 0x0B)),
    ("Stockage", "MongoDB +\nMetadonnees", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (title, desc, color) in enumerate(pipeline_steps):
    x = Inches(0.6 + i * 2.5)
    y = Inches(2.9)

    add_rounded_rect(slide, x, y, Inches(2.1), Inches(1.5), LIGHT_NAVY, name=f"s2_pipe_{i}")
    add_circle(slide, x + Inches(0.8), y + Inches(0.15), Inches(0.5), color)
    add_text(slide, x + Inches(0.8), y + Inches(0.22), Inches(0.5), Inches(0.35),
             f"{i+1}", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.75), Inches(2.1), Inches(0.3),
             title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.05), Inches(2.1), Inches(0.4),
             desc.replace('\n', ' '), font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(pipeline_steps) - 1:
        add_text(slide, x + Inches(2.0), y + Inches(0.5), Inches(0.5), Inches(0.4),
                 "\u2192", font_size=20, color=EMERALD, alignment=PP_ALIGN.CENTER)

# Extracted fields
add_text(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
         "Champs Extraits Automatiquement", font_size=16, color=EMERALD, bold=True)

fields = [
    ["Vendeur", "N\u00b0 Facture", "Date", "Echeance"],
    ["TVA Vendeur", "TVA Acheteur", "Lignes", "Sous-total"],
    ["Taxe", "Remise", "Total TTC", "Devise"],
]

for row_i, row in enumerate(fields):
    for col_i, field in enumerate(row):
        x = Inches(0.8 + col_i * 3.0)
        y = Inches(5.3 + row_i * 0.6)
        add_rounded_rect(slide, x, y, Inches(2.7), Inches(0.45), LIGHT_NAVY)
        add_text(slide, x + Inches(0.15), y + Inches(0.07), Inches(2.4), Inches(0.3),
                 f"\u2713  {field}", font_size=11, color=GRAY)

add_footer(slide, 11)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Sprint 2 : Verification Double Ecran
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), EMERALD, name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 2", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "Verification & Scoring de Confiance", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Split screen mockup
# Left panel: Document
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), LIGHT_NAVY, name="verify_left")
add_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.5), BLUE)
add_text(slide, Inches(1.0), Inches(2.35), Inches(5), Inches(0.4),
         "Document Original", font_size=13, color=WHITE, bold=True)

# Mock document content
add_rounded_rect(slide, Inches(1.2), Inches(3.0), Inches(4.7), Inches(3.5),
                 RGBColor(0x14, 0x38, 0x5C))
doc_lines = [
    "FACTURE N\u00b0 INV-2026-0042",
    "Date: 15/03/2026",
    "Vendeur: TechSoft SARL",
    "TVA: 1234567A/P/M/000",
    "",
    "Designation     Qte   PU      Total",
    "Service Dev     1     5000    5000",
    "Maintenance     1     1200    1200",
    "",
    "Sous-total:           6200 TND",
    "TVA 19%:              1178 TND",
    "Total TTC:            7378 TND",
]
for j, line in enumerate(doc_lines):
    add_text(slide, Inches(1.5), Inches(3.1 + j * 0.27), Inches(4.2), Inches(0.25),
             line, font_size=9, color=GRAY, font_name="Consolas")

# Right panel: Extracted data
add_rounded_rect(slide, Inches(6.6), Inches(2.3), Inches(6.0), Inches(4.5), LIGHT_NAVY, name="verify_right")
add_rect(slide, Inches(6.6), Inches(2.3), Inches(6.0), Inches(0.5), EMERALD)
add_text(slide, Inches(6.8), Inches(2.35), Inches(5), Inches(0.4),
         "Donnees Extraites + Confiance", font_size=13, color=WHITE, bold=True)

# Confidence bars
conf_fields = [
    ("Vendeur: TechSoft SARL", 95, EMERALD),
    ("N\u00b0 Facture: INV-2026-0042", 98, EMERALD),
    ("Date: 2026-03-15", 92, EMERALD),
    ("TVA Vendeur: 1234567A/P/M/000", 88, EMERALD),
    ("Sous-total: 6200 TND", 90, EMERALD),
    ("TVA: 1178 TND (19%)", 85, RGBColor(0xF5, 0x9E, 0x0B)),
    ("Total TTC: 7378 TND", 95, EMERALD),
]

for j, (field, conf, color) in enumerate(conf_fields):
    fy = Inches(3.0 + j * 0.5)
    add_text(slide, Inches(6.9), fy, Inches(3.5), Inches(0.25),
             field, font_size=10, color=WHITE)

    # Confidence bar background
    add_rounded_rect(slide, Inches(10.5), fy + Inches(0.05), Inches(1.5), Inches(0.18),
                     RGBColor(0x1A, 0x3A, 0x55))
    # Confidence bar fill
    bar_width = 1.5 * (conf / 100)
    add_rounded_rect(slide, Inches(10.5), fy + Inches(0.05), Inches(bar_width), Inches(0.18), color)
    # Percentage
    add_text(slide, Inches(12.05), fy, Inches(0.5), Inches(0.25),
             f"{conf}%", font_size=9, color=color, bold=True)

# Review threshold note
add_rounded_rect(slide, Inches(6.9), Inches(6.2), Inches(5.4), Inches(0.4), RGBColor(0x14, 0x38, 0x5C))
add_text(slide, Inches(7.1), Inches(6.25), Inches(5), Inches(0.3),
         "Seuil de review humain : < 60% confiance sur champs critiques",
         font_size=10, color=RGBColor(0xF5, 0x9E, 0x0B))

add_footer(slide, 12)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — Sprint 2 : Metriques
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), EMERALD, name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 2", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "Resultats Sprint 2", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

metrics_s2 = [
    ("557", "Lignes Extraction IA", BLUE),
    ("3", "Modeles LLM Integres", EMERALD),
    ("50+", "Devises Supportees", RGBColor(0xF5, 0x9E, 0x0B)),
    ("12", "Champs Extraits", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (val, label, color) in enumerate(metrics_s2):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.5)
    add_metric_card(slide, x, y, Inches(2.8), Inches(1.6), val, label, color, name=f"s2_metric_{i}")

# Deliverables
add_text(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.4),
         "Livrables Sprint 2", font_size=18, color=EMERALD, bold=True)

deliverables_s2 = [
    "Pipeline triple fallback : Claude -> Ollama -> Template",
    "OCR multi-langue : Francais, Anglais, Arabe (RTL)",
    "Scoring de confiance 0-100 par champ extrait",
    "Interface de verification double ecran",
    "Traitement batch asynchrone via Celery",
    "Detection automatique de 50+ devises",
    "Validation arithmetique intelligente",
]

for i, d in enumerate(deliverables_s2):
    y = Inches(5.0 + i * 0.35)
    add_text(slide, Inches(1.0), y, Inches(11), Inches(0.3),
             f"\u2713  {d}", font_size=12, color=GRAY)

add_footer(slide, 13)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Sprint 3 : El Fatoora & Billing
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), RGBColor(0xF5, 0x9E, 0x0B), name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 3", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "El Fatoora, Signature & Billing", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# El Fatoora submission flow
add_text(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.4),
         "Flux de Soumission El Fatoora", font_size=16, color=EMERALD, bold=True)

ef_steps = [
    ("Pre-Validation", "Champs obligatoires\nMatricule fiscale", BLUE),
    ("Validation Fiscale", "TVA, arithmetique\nScore conformite", EMERALD),
    ("Generation XML", "Format TTN\nStructure normee", RGBColor(0xF5, 0x9E, 0x0B)),
    ("Soumission", "API El Fatoora\nID unique TTN", RGBColor(0xEF, 0x44, 0x44)),
    ("Audit Trail", "Horodatage\nHistorique complet", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(ef_steps):
    x = Inches(0.5 + i * 2.5)
    y = Inches(2.8)

    add_rounded_rect(slide, x, y, Inches(2.15), Inches(1.4), LIGHT_NAVY, name=f"ef_step_{i}")
    add_rect(slide, x, y, Inches(2.15), Pt(3), color)
    add_text(slide, x, y + Inches(0.15), Inches(2.15), Inches(0.3),
             title, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.5), Inches(2.15), Inches(0.7),
             desc.replace('\n', ' '), font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(ef_steps) - 1:
        add_text(slide, x + Inches(2.05), y + Inches(0.35), Inches(0.5), Inches(0.4),
                 "\u2192", font_size=18, color=EMERALD, alignment=PP_ALIGN.CENTER)

# Billing plans
add_text(slide, Inches(0.8), Inches(4.6), Inches(5), Inches(0.4),
         "Plans Tarifaires Stripe", font_size=16, color=EMERALD, bold=True)

plans = [
    ("Free", "0 TND/mois", "50 factures, 3 membres, 20 requetes IA", GRAY),
    ("Pro", "29 USD/mois", "500 factures, 15 membres, 200 requetes IA", BLUE),
    ("Enterprise", "99 USD/mois", "Illimite, support prioritaire, SLA", EMERALD),
]

for i, (plan, price, desc, color) in enumerate(plans):
    x = Inches(0.8 + i * 4.0)
    y = Inches(5.1)

    add_rounded_rect(slide, x, y, Inches(3.6), Inches(1.8), LIGHT_NAVY, name=f"plan_{i}")
    add_text(slide, x + Inches(0.3), y + Inches(0.15), Inches(1.5), Inches(0.35),
             plan, font_size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.55), Inches(3), Inches(0.35),
             price, font_size=22, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), Inches(3), Inches(0.6),
             desc, font_size=11, color=GRAY)

add_footer(slide, 14)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — Sprint 3 : Dashboard & Conseiller IA
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), RGBColor(0xF5, 0x9E, 0x0B), name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 3", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "Dashboard & Conseiller IA", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Dashboard mockup
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(7.5), Inches(4.8), LIGHT_NAVY, name="dash_mock")
add_rect(slide, Inches(0.8), Inches(2.3), Inches(7.5), Inches(0.5), BLUE)
add_text(slide, Inches(1.0), Inches(2.35), Inches(6), Inches(0.4),
         "Tableau de Bord Analytique", font_size=14, color=WHITE, bold=True)

# KPI mockup cards
kpi_data = [
    ("156 420", "Montant Total", EMERALD),
    ("342", "Factures", BLUE),
    ("12", "A Revoir", RGBColor(0xF5, 0x9E, 0x0B)),
    ("457", "Moy/Facture", ACCENT_BLUE),
]
for i, (val, label, color) in enumerate(kpi_data):
    x = Inches(1.1 + i * 1.75)
    add_rounded_rect(slide, x, Inches(3.0), Inches(1.55), Inches(0.9), RGBColor(0x14, 0x38, 0x5C))
    add_text(slide, x + Inches(0.1), Inches(3.05), Inches(1.35), Inches(0.35),
             val, font_size=14, color=color, bold=True)
    add_text(slide, x + Inches(0.1), Inches(3.4), Inches(1.35), Inches(0.3),
             label, font_size=9, color=GRAY)

# Chart area mockup
add_rounded_rect(slide, Inches(1.1), Inches(4.1), Inches(4.0), Inches(2.7), RGBColor(0x14, 0x38, 0x5C))
add_text(slide, Inches(1.3), Inches(4.2), Inches(3.5), Inches(0.3),
         "Tendance Mensuelle", font_size=11, color=EMERALD, bold=True)
# Simulated chart bars
for j in range(6):
    bar_h = [0.8, 1.2, 0.9, 1.5, 1.1, 1.8][j]
    add_rect(slide, Inches(1.6 + j * 0.55), Inches(6.5 - bar_h), Inches(0.35), Inches(bar_h), BLUE)

# Pie chart area
add_rounded_rect(slide, Inches(5.3), Inches(4.1), Inches(2.7), Inches(2.7), RGBColor(0x14, 0x38, 0x5C))
add_text(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.3),
         "Par Vendeur", font_size=11, color=EMERALD, bold=True)
add_circle(slide, Inches(5.9), Inches(4.7), Inches(1.5), BLUE)
add_circle(slide, Inches(6.15), Inches(4.95), Inches(1.0), EMERALD)
add_circle(slide, Inches(6.35), Inches(5.15), Inches(0.6), RGBColor(0xF5, 0x9E, 0x0B))

# AI Advisor panel
add_rounded_rect(slide, Inches(8.6), Inches(2.3), Inches(4.2), Inches(4.8), LIGHT_NAVY, name="advisor_mock")
add_rect(slide, Inches(8.6), Inches(2.3), Inches(4.2), Inches(0.5), EMERALD)
add_text(slide, Inches(8.8), Inches(2.35), Inches(3.8), Inches(0.4),
         "Conseiller IA", font_size=14, color=WHITE, bold=True)

# Chat messages
chat = [
    ("Utilisateur", "Analyser mes depenses Q1 2026", False),
    ("IA", "Vos depenses Q1 montrent une hausse de 15% vs Q4. Principaux postes : services IT (42%), fournitures (28%).", True),
    ("Utilisateur", "Quelles factures sont non-conformes ?", False),
    ("IA", "3 factures n'ont pas de matricule fiscale vendeur. Recommandation : contacter les fournisseurs.", True),
]

cy = Inches(3.0)
for sender, msg, is_ai in chat:
    bg_c = RGBColor(0x14, 0x38, 0x5C) if is_ai else RGBColor(0x1A, 0x4A, 0x6C)
    msg_w = Inches(3.5)
    msg_x = Inches(8.8) if is_ai else Inches(9.1)
    add_rounded_rect(slide, msg_x, cy, msg_w, Inches(0.7), bg_c)
    add_text(slide, msg_x + Inches(0.1), cy + Inches(0.05), msg_w - Inches(0.2), Inches(0.15),
             sender, font_size=8, color=EMERALD if is_ai else BLUE, bold=True)
    add_text(slide, msg_x + Inches(0.1), cy + Inches(0.22), msg_w - Inches(0.2), Inches(0.45),
             msg, font_size=8, color=GRAY)
    cy += Inches(0.85)

add_footer(slide, 15)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — Sprint 3 : Metriques
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rounded_rect(slide, Inches(0.8), Inches(0.5), Inches(1.8), Inches(0.5), RGBColor(0xF5, 0x9E, 0x0B), name="sprint_badge")
add_text(slide, Inches(0.8), Inches(0.55), Inches(1.8), Inches(0.4),
         "SPRINT 3", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.7),
         "Resultats Sprint 3", font_size=36, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

metrics_s3 = [
    ("1 043", "Lignes TEJ XML", BLUE),
    ("9", "Routes TTN", EMERALD),
    ("3", "Plans Stripe", RGBColor(0xF5, 0x9E, 0x0B)),
    ("1 209", "Lignes Notifications", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (val, label, color) in enumerate(metrics_s3):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.5)
    add_metric_card(slide, x, y, Inches(2.8), Inches(1.6), val, label, color, name=f"s3_metric_{i}")

add_text(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.4),
         "Livrables Sprint 3", font_size=18, color=EMERALD, bold=True)

deliverables_s3 = [
    "Integration El Fatoora TTN avec soumission et audit trail",
    "TEJ : certificats de retenue a la source (5 types)",
    "Signature electronique XAdES-B avec validation",
    "Tableau de bord analytique avec 4 KPIs + 3 graphiques",
    "Conseiller IA conversationnel integre",
    "Systeme de notifications multi-canal (in-app, email, SMS)",
    "Integration Stripe avec 3 plans tarifaires",
    "Export Excel individuel et batch",
]

for i, d in enumerate(deliverables_s3):
    y = Inches(4.95 + i * 0.32)
    add_text(slide, Inches(1.0), y, Inches(11), Inches(0.3),
             f"\u2713  {d}", font_size=12, color=GRAY)

add_footer(slide, 16)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17 — Fonctionnalites Avancees
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.7),
         "Fonctionnalites Avancees", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

features_adv = [
    ("QR Code Partage", "Generation et scan de QR codes\npour partage inter-utilisateurs", BLUE),
    ("Trilingue FR/EN/AR", "Support complet avec\ngestion RTL arabe", EMERALD),
    ("Notifications", "In-app, Email (Gmail SMTP),\nSMS (Twilio) + Digest", RGBColor(0xF5, 0x9E, 0x0B)),
    ("Audit Trail", "Historique complet des actions\navec IP et horodatage", RGBColor(0xEF, 0x44, 0x44)),
    ("Export Multi-Format", "Excel (openpyxl), PDF (jsPDF)\navec mise en forme", ACCENT_BLUE),
    ("Gestion d'Equipe", "Invitations, roles, permissions\npar organisation", SOFT_GREEN),
]

for i, (title, desc, color) in enumerate(features_adv):
    x = Inches(0.8 + (i % 3) * 4.0)
    y = Inches(2.0 + (i // 3) * 2.5)

    add_rounded_rect(slide, x, y, Inches(3.6), Inches(2.1), LIGHT_NAVY, name=f"adv_feat_{i}")
    add_rect(slide, x, y, Pt(4), Inches(2.1), color)
    add_text(slide, x + Inches(0.25), y + Inches(0.2), Inches(3.1), Inches(0.4),
             title, font_size=18, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.7), Inches(3.1), Inches(1.1),
             desc.replace('\n', ' '), font_size=13, color=GRAY)

add_footer(slide, 17)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18 — Metriques Globales
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.7),
         "Metriques Globales du Projet", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Big metrics grid
big_metrics = [
    ("35 996", "Lignes de Code", BLUE),
    ("51+", "Endpoints API", EMERALD),
    ("16", "Microservices Backend", RGBColor(0xF5, 0x9E, 0x0B)),
    ("68+", "Composants React", RGBColor(0xEF, 0x44, 0x44)),
    ("11", "Collections MongoDB", ACCENT_BLUE),
    ("3", "Langues (FR/EN/AR)", SOFT_GREEN),
    ("17", "Pages Frontend", BLUE),
    ("130", "Fichiers Source", EMERALD),
]

for i, (val, label, color) in enumerate(big_metrics):
    x = Inches(0.5 + (i % 4) * 3.1)
    y = Inches(2.0 + (i // 4) * 2.5)
    add_metric_card(slide, x, y, Inches(2.8), Inches(2.0), val, label, color, name=f"global_metric_{i}")

add_footer(slide, 18)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 19 — Code Breakdown
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.7),
         "Repartition du Code", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

# Left: Backend breakdown
add_text(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(0.4),
         "Backend Python — 10 205 lignes", font_size=18, color=BLUE, bold=True)

backend_modules = [
    ("app.py (API Routes)", "2 671", 26),
    ("notification_service.py", "1 209", 12),
    ("tej_xml_generator.py", "1 043", 10),
    ("org_service.py", "785", 8),
    ("extractor_core.py", "557", 5),
    ("xml_parser.py", "510", 5),
    ("tax_validator.py", "391", 4),
    ("billing_service.py", "358", 3),
    ("Autres services (8)", "2 681", 27),
]

for i, (module, lines, pct) in enumerate(backend_modules):
    y = Inches(2.4 + i * 0.52)
    add_text(slide, Inches(0.8), y, Inches(3.5), Inches(0.3),
             module, font_size=11, color=GRAY)
    add_text(slide, Inches(4.3), y, Inches(0.8), Inches(0.3),
             lines, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    # Bar
    bar_w = 1.0 * (pct / 27)
    add_rounded_rect(slide, Inches(5.3), y + Inches(0.05), Inches(bar_w), Inches(0.18), BLUE)

# Right: Frontend breakdown
add_text(slide, Inches(7.0), Inches(1.9), Inches(5), Inches(0.4),
         "Frontend TypeScript — 25 791 lignes", font_size=18, color=EMERALD, bold=True)

frontend_modules = [
    ("Composants Metier (22)", "~12 000", 46),
    ("Composants UI Shadcn (46)", "~8 000", 31),
    ("Pages Next.js (17)", "~3 000", 12),
    ("Contexts & Lib", "~1 500", 6),
    ("Configuration", "~1 291", 5),
]

for i, (module, lines, pct) in enumerate(frontend_modules):
    y = Inches(2.4 + i * 0.7)
    add_text(slide, Inches(7.0), y, Inches(3.5), Inches(0.3),
             module, font_size=12, color=GRAY)
    add_text(slide, Inches(10.5), y, Inches(1.0), Inches(0.3),
             lines, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    # Bar
    bar_w = 1.5 * (pct / 46)
    add_rounded_rect(slide, Inches(7.0), y + Inches(0.35), Inches(bar_w), Inches(0.2), EMERALD)

# Ratio
add_rounded_rect(slide, Inches(3.0), Inches(6.2), Inches(7.3), Inches(0.8), LIGHT_NAVY)
add_text(slide, Inches(3.0), Inches(6.3), Inches(7.3), Inches(0.5),
         "Ratio Frontend/Backend :  72% / 28%  |  130 fichiers source  |  5 services Docker",
         font_size=13, color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

add_footer(slide, 19)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 20 — Difficultes & Solutions
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.7),
         "Difficultes & Solutions", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

challenges = [
    ("Precision OCR Arabe",
     "Tesseract peu fiable sur texte arabe manuscrit",
     "Pre-traitement image (contraste, debruitage) + prompt LLM specialise RTL",
     RGBColor(0xEF, 0x44, 0x44)),
    ("Disponibilite LLM",
     "API Cloud indisponible = blocage extraction",
     "Architecture triple fallback : Claude -> Ollama local -> Template",
     BLUE),
    ("Conformite El Fatoora",
     "API TTN non encore publiee officiellement",
     "Simulation complete + structure preparee pour integration reelle",
     RGBColor(0xF5, 0x9E, 0x0B)),
    ("Isolation Multi-Tenant",
     "Risque de fuite de donnees entre organisations",
     "Filtrage systematique par org_id + decorateurs d'acces",
     EMERALD),
]

for i, (title, problem, solution, color) in enumerate(challenges):
    y = Inches(1.9 + i * 1.35)

    add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.15), LIGHT_NAVY, name=f"challenge_{i}")
    add_rect(slide, Inches(0.8), y, Pt(4), Inches(1.15), color)

    add_text(slide, Inches(1.1), y + Inches(0.1), Inches(3), Inches(0.3),
             title, font_size=16, color=WHITE, bold=True)

    add_text(slide, Inches(1.1), y + Inches(0.45), Inches(5), Inches(0.3),
             f"Probleme : {problem}", font_size=11, color=RGBColor(0xEF, 0x44, 0x44))

    add_text(slide, Inches(1.1), y + Inches(0.75), Inches(10.8), Inches(0.3),
             f"Solution : {solution}", font_size=11, color=EMERALD)

add_footer(slide, 20)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 21 — Perspectives d'Evolution
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

add_circle(slide, Inches(-2), Inches(-2), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(11), Inches(5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

add_rect(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4), EMERALD, name="accent_line")
add_text(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.7),
         "Perspectives d'Evolution", font_size=38, color=WHITE, bold=True,
         font_name="Segoe UI Light", name="section_title")

perspectives = [
    ("Court Terme", BLUE, [
        "Integration API El Fatoora reelle",
        "Application mobile React Native",
        "Amelioration OCR avec modeles dedies",
    ]),
    ("Moyen Terme", EMERALD, [
        "Marketplace de templates facture",
        "Integration comptable (Sage, QuickBooks)",
        "IA predictive pour tresorerie",
    ]),
    ("Long Terme", RGBColor(0xF5, 0x9E, 0x0B), [
        "Expansion multi-pays (Maghreb, Afrique)",
        "Blockchain pour audit immutable",
        "API ouverte pour ecosysteme tiers",
    ]),
]

for i, (title, color, items) in enumerate(perspectives):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.0)

    add_rounded_rect(slide, x, y, Inches(3.6), Inches(4.5), LIGHT_NAVY, name=f"persp_{i}")
    add_rect(slide, x, y, Inches(3.6), Inches(0.6), color)
    add_text(slide, x, y + Inches(0.1), Inches(3.6), Inches(0.45),
             title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        iy = y + Inches(0.9 + j * 1.1)
        add_rounded_rect(slide, x + Inches(0.2), iy, Inches(3.2), Inches(0.85),
                         RGBColor(0x14, 0x38, 0x5C))
        add_text(slide, x + Inches(0.4), iy + Inches(0.2), Inches(2.8), Inches(0.5),
                 item, font_size=13, color=GRAY)

add_footer(slide, 21)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 22 — Conclusion / Merci
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
set_morph(slide)

# Decorative elements
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_1")
add_circle(slide, Inches(10), Inches(4.5), Inches(5), RGBColor(0x0E, 0x2F, 0x55), name="deco_circle_2")

# Center everything
add_rect(slide, Inches(5.6), Inches(1.8), Inches(2.1), Pt(4), EMERALD, name="accent_line")

add_text(slide, Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.2),
         "Merci de votre attention", font_size=48, color=WHITE, bold=True,
         font_name="Segoe UI Light", alignment=PP_ALIGN.CENTER, name="main_title")

add_text(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.8),
         "SmartInvoice AI", font_size=32, color=EMERALD,
         font_name="Segoe UI Light", alignment=PP_ALIGN.CENTER, name="tagline")

add_text(slide, Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.6),
         "Plateforme SaaS de Gestion Intelligente des Factures", font_size=18,
         color=GRAY, font_name="Segoe UI Light", alignment=PP_ALIGN.CENTER, name="tagline2")

# Key metrics summary
summary_metrics = [
    ("35 996", "Lignes de Code"),
    ("51+", "API Endpoints"),
    ("16", "Microservices"),
    ("3", "Sprints"),
]

for i, (val, label) in enumerate(summary_metrics):
    x = Inches(2.0 + i * 2.5)
    y = Inches(5.0)
    add_text(slide, x, y, Inches(2.2), Inches(0.5),
             val, font_size=28, color=EMERALD, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")
    add_text(slide, x, y + Inches(0.5), Inches(2.2), Inches(0.3),
             label, font_size=12, color=GRAY,
             alignment=PP_ALIGN.CENTER)

# Student info
add_text(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.4),
         "Anas Badri  |  2025-2026  |  github.com/myanasbadri/PFE_SAAS",
         font_size=13, color=GRAY, font_name="Segoe UI Light",
         alignment=PP_ALIGN.CENTER, name="student_info")

# Dots
for i in range(5):
    add_circle(slide, Inches(5.8 + i * 0.25), Inches(6.7), Inches(0.08),
               EMERALD if i == 2 else RGBColor(0x1A, 0x3A, 0x55), name=f"dot_{i}")


# ═══════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════
output_path = "C:/Users/MSI/Desktop/PFE/SmartInvoice_AI_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
